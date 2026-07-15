"""
apps/documents/pdf_tools.py

Stateless PDF operations for the in-app PDF editor's server-side jobs
(compression, conversion). The editor handles everything it can in the browser
and POSTs only the heavy jobs here — see PdfToolView.

Reuses the project's existing tooling:
  • PyMuPDF (fitz) for text extraction and lightweight re-compression
  • Ghostscript for deep compression when available
  • LibreOffice for PDF → Office conversions (same binary as the preview pipeline)
"""
from __future__ import annotations

import io
import logging
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from django.conf import settings as django_settings

logger = logging.getLogger(__name__)


# ── LibreOffice / Ghostscript discovery ──────────────────────────────────────

def resolve_soffice_bin() -> str | None:
    """Locate the LibreOffice binary, mirroring generate_document_preview."""
    configured = (
        getattr(django_settings, "LIBREOFFICE_BIN", "").strip()
        or getattr(django_settings, "LIBREOFFICE_CMD", "").strip()
    )
    if configured and Path(configured).exists():
        return configured
    return (
        shutil.which("libreoffice")
        or shutil.which("soffice")
        or ("/usr/bin/libreoffice" if Path("/usr/bin/libreoffice").exists() else None)
        or ("/usr/bin/soffice" if Path("/usr/bin/soffice").exists() else None)
    )


def _ghostscript_bin() -> str | None:
    return shutil.which("gs") or shutil.which("gswin64c") or shutil.which("gswin32c")


# ── Compression ───────────────────────────────────────────────────────────────

# Map the editor's levels to Ghostscript's quality presets.
_GS_PDFSETTINGS = {"high": "/ebook", "extreme": "/screen"}


def compress_pdf(data: bytes, level: str = "high") -> bytes:
    """
    Reduce a PDF's size. Prefers Ghostscript (which downsamples images — the bulk
    of most documents); falls back to a PyMuPDF re-save (deflate + garbage
    collection) when Ghostscript isn't installed. Returns the original bytes if
    nothing managed to shrink it, so callers never get a *larger* file.
    """
    gs = _ghostscript_bin()
    if gs:
        setting = _GS_PDFSETTINGS.get(level, "/ebook")
        with tempfile.TemporaryDirectory(prefix="pdfcompress_") as d:
            src = Path(d) / "in.pdf"
            out = Path(d) / "out.pdf"
            src.write_bytes(data)
            cmd = [
                gs, "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
                f"-dPDFSETTINGS={setting}", "-dNOPAUSE", "-dQUIET", "-dBATCH",
                "-dDetectDuplicateImages=true", "-dCompressFonts=true",
                f"-sOutputFile={out}", str(src),
            ]
            try:
                subprocess.run(
                    cmd, timeout=180, check=True,
                    stdout=subprocess.DEVNULL, stderr=subprocess.PIPE,
                )
                if out.exists():
                    result = out.read_bytes()
                    if result and len(result) < len(data):
                        return result
            except Exception as exc:
                logger.warning("Ghostscript compression failed, falling back: %s", exc)

    fitted = _fitz_compress(data)
    return fitted if len(fitted) < len(data) else data


def _fitz_compress(data: bytes) -> bytes:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        return doc.tobytes(
            garbage=4, deflate=True, clean=True,
            deflate_images=True, deflate_fonts=True,
        )
    finally:
        doc.close()


# ── Redaction ─────────────────────────────────────────────────────────────────

def redact_pdf(data: bytes, rects: list[dict]) -> bytes:
    """
    Permanently remove content under each rectangle (true redaction, not a cosmetic
    black box): PyMuPDF's apply_redactions rewrites the page content stream, dropping
    text, vector and image data that intersects the area, then paints it black.

    `rects` items are page fractions: {page, x, y, w, h} with x/y measured from the
    page's top-left, matching the editor's annotation coordinates.
    """
    import fitz  # PyMuPDF

    by_page: dict[int, list[dict]] = {}
    for r in rects:
        try:
            by_page.setdefault(int(r["page"]), []).append(r)
        except (KeyError, TypeError, ValueError):
            continue

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for pno, items in by_page.items():
            if pno < 0 or pno >= doc.page_count:
                continue
            page = doc[pno]
            pw, ph = page.rect.width, page.rect.height
            for r in items:
                try:
                    x0 = float(r["x"]) * pw
                    y0 = float(r["y"]) * ph
                    x1 = x0 + float(r["w"]) * pw
                    y1 = y0 + float(r["h"]) * ph
                except (KeyError, TypeError, ValueError):
                    continue
                page.add_redact_annot(fitz.Rect(x0, y0, x1, y1), fill=(0, 0, 0))
            page.apply_redactions()
        return doc.tobytes(garbage=4, deflate=True, clean=True)
    finally:
        doc.close()


# ── Text extraction ───────────────────────────────────────────────────────────

def pdf_to_text_bytes(data: bytes) -> bytes:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        text = "\n\n".join(page.get_text() for page in doc)
    finally:
        doc.close()
    return text.encode("utf-8")


# ── LibreOffice conversion (PDF → Office) ─────────────────────────────────────

def libreoffice_convert(data: bytes, src_suffix: str, convert_to: str, timeout: int = 180) -> bytes:
    """
    Convert a document to another format with a headless, isolated LibreOffice
    process. Used for PDF → docx/xlsx/pptx (best-effort: LibreOffice imports the
    PDF via its Draw filter, so fidelity varies by document).
    """
    soffice = resolve_soffice_bin()
    if not soffice:
        raise RuntimeError("Document conversion is unavailable: LibreOffice is not installed on the server.")

    # Reuse the preview pipeline's robust child-process cleanup.
    from .tasks import _kill_process_tree

    with tempfile.TemporaryDirectory(prefix="pdftool_") as d:
        tmp = Path(d)
        profile = tmp / "profile"
        (profile / "data").mkdir(parents=True)
        (profile / "cache").mkdir(parents=True)
        outdir = tmp / "out"
        outdir.mkdir()
        src = tmp / f"in{src_suffix}"
        src.write_bytes(data)

        env = os.environ.copy()
        env.update({
            "DISPLAY": ":99",
            "HOME": str(tmp),
            "XDG_CONFIG_HOME": str(profile),
            "XDG_DATA_HOME": str(profile / "data"),
            "XDG_CACHE_HOME": str(profile / "cache"),
        })

        cmd = [
            soffice,
            f"-env:UserInstallation={profile.as_uri()}",
            "--headless", "--nocrashreport", "--nodefault",
            "--nofirststartwizard", "--nologo", "--norestore",
            "--convert-to", convert_to,
            "--outdir", str(outdir),
            str(src),
        ]
        popen_kwargs = (
            {"creationflags": subprocess.CREATE_NEW_PROCESS_GROUP}
            if os.name == "nt" else {"start_new_session": True}
        )
        proc = subprocess.Popen(
            cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, env=env, **popen_kwargs,
        )
        try:
            out, _ = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            _kill_process_tree(proc)
            try:
                proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                _kill_process_tree(proc, force=True)
                proc.wait()
            raise RuntimeError("Conversion timed out.")

        produced = [p for p in sorted(outdir.glob("*")) if p.is_file()]
        if proc.returncode != 0 or not produced:
            tail = (out or b"").decode("utf-8", "replace")[:500]
            logger.error("LibreOffice convert-to %s failed (exit=%s): %s", convert_to, proc.returncode, tail)
            raise RuntimeError("Could not convert this document — the format may not be supported.")
        return produced[0].read_bytes()


# Source types LibreOffice can turn into PDF (covers the editor's "open Office
# document" flow). Anything else falls back to a .docx assumption.
_OFFICE_SUFFIXES = {
    ".doc", ".docx", ".docm", ".dot", ".dotx", ".rtf", ".odt", ".txt",
    ".xls", ".xlsx", ".xlsm", ".xlsb", ".ods", ".csv",
    ".ppt", ".pptx", ".pptm", ".ppsx", ".odp",
    ".html", ".htm",
}


def office_to_pdf(data: bytes, filename: str) -> bytes:
    """Convert an Office/HTML document (identified by filename suffix) to PDF."""
    suffix = os.path.splitext(filename or "")[1].lower()
    if suffix not in _OFFICE_SUFFIXES:
        suffix = suffix or ".docx"
    return libreoffice_convert(data, suffix, "pdf")


# ── PDF → Office ──────────────────────────────────────────────────────────────
# LibreOffice can't export a PDF (which it opens in Draw) to docx/xlsx/pptx, so
# these use purpose-built libraries instead.

def pdf_to_docx(data: bytes) -> bytes:
    """High-fidelity PDF → Word using pdf2docx (preserves layout, tables, images)."""
    try:
        from pdf2docx import Converter
    except ImportError:
        raise RuntimeError("PDF → Word isn't available on the server (pdf2docx not installed).")
    with tempfile.TemporaryDirectory(prefix="pdf2docx_") as d:
        src = Path(d) / "in.pdf"
        out = Path(d) / "out.docx"
        src.write_bytes(data)
        cv = Converter(str(src))
        try:
            cv.convert(str(out))
        finally:
            cv.close()
        if not out.exists():
            raise RuntimeError("Could not convert this PDF to Word.")
        return out.read_bytes()


def pdf_to_pptx(data: bytes) -> bytes:
    """PDF → PowerPoint: one slide per page, each page rendered as a full-bleed image."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        raise RuntimeError("PDF → PowerPoint isn't available on the server (python-pptx not installed).")
    import fitz

    EMU_PER_PT = 12700
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        if doc.page_count:
            r = doc[0].rect
            prs.slide_width = Emu(int(r.width * EMU_PER_PT))
            prs.slide_height = Emu(int(r.height * EMU_PER_PT))
        zoom = fitz.Matrix(150 / 72, 150 / 72)  # ~150 DPI
        for page in doc:
            pix = page.get_pixmap(matrix=zoom)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
    finally:
        doc.close()
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def pdf_to_xlsx(data: bytes) -> bytes:
    """
    Best-effort PDF → Excel: detected tables become sheets; pages without tables
    fall back to one text line per row. (PDFs aren't spreadsheets, so fidelity
    depends heavily on the source.)
    """
    try:
        from openpyxl import Workbook
    except ImportError:
        raise RuntimeError("PDF → Excel isn't available on the server (openpyxl not installed).")
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    wb = Workbook()
    wb.remove(wb.active)
    try:
        for pno, page in enumerate(doc, start=1):
            try:
                found = page.find_tables()
                tables = list(found.tables) if found else []
            except Exception:
                tables = []
            if tables:
                for ti, tbl in enumerate(tables, start=1):
                    ws = wb.create_sheet(title=f"P{pno}-T{ti}"[:31])
                    for row in tbl.extract():
                        ws.append([("" if c is None else str(c)) for c in row])
            else:
                ws = wb.create_sheet(title=f"Page {pno}"[:31])
                for line in page.get_text().splitlines():
                    ws.append([line])
    finally:
        doc.close()
    if not wb.sheetnames:
        wb.create_sheet(title="Sheet1")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
