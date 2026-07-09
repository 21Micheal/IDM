from io import BytesIO
from pathlib import Path
import tempfile
import os
import re
import logging
from decimal import Decimal, InvalidOperation

from apps.search.utils import SEARCH_INDEX_EXCEPTIONS

logger = logging.getLogger(__name__)


STANDARD_DOCUMENT_FIELDS = {"title", "supplier", "amount", "currency", "document_date", "due_date"}


def _stringish(value):
    if value is None:
        return ""
    if isinstance(value, (str, int, float, Decimal)):
        return str(value)
    return ""


def _form_values_for_metadata(values: dict) -> dict:
    metadata = {}
    for key, value in (values or {}).items():
        if key in STANDARD_DOCUMENT_FIELDS:
            continue
        if isinstance(value, dict) and value.get("storage_path"):
            continue
        metadata[key] = value
    return metadata


def _document_field_kwargs(values: dict) -> dict:
    fields = {}
    for key in ("supplier", "currency", "document_date", "due_date"):
        value = _stringish((values or {}).get(key)).strip()
        if value:
            fields[key] = value
    amount = _stringish((values or {}).get("amount")).replace(",", "").strip()
    if amount:
        try:
            fields["amount"] = Decimal(amount)
        except (InvalidOperation, ValueError):
            pass
    return fields


def _display_value(value):
    if isinstance(value, dict) and value.get("storage_path"):
        return value.get("name") or "Attached file"
    if isinstance(value, (list, dict)):
        return str(value)
    return value


def _decode_data_url_image(value):
    """Decode a ``data:image/...;base64,...`` URL (e.g. a signature) to raw
    bytes, or return None. Prevents dumping a giant base64 string into the doc."""
    if not isinstance(value, str) or not value.startswith("data:image"):
        return None
    try:
        import base64
        return base64.b64decode(value.split(",", 1)[1])
    except Exception:
        return None


# ─── Helpers ────────────────────────────────────────────────────────────────

def _replace_placeholder_in_paragraph(para, values: dict):
    """
    Replace {{key}} placeholders in a paragraph, handling the case where
    placeholders are split across multiple runs (common in Word).
    We stitch the full text, replace, then rewrite into the first run.
    """
    full_text = "".join(run.text for run in para.runs)
    if "{{" not in full_text:
        return

    def replacer(match):
        key = match.group(1)
        val = values.get(key)
        if val is None:
            return ""
        if isinstance(val, (list, dict)):
            return str(val)
        return str(val)

    new_text = re.sub(r"\{\{([a-zA-Z0-9_]+)\}\}", replacer, full_text)
    if new_text == full_text:
        return

    # Write into first run, clear the rest
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""


# ─── Built template generators ───────────────────────────────────────────────

def generate_built_pdf(template, values) -> bytes:
    """Generate PDF from built template using reportlab."""
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable,
    )
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=20*mm, bottomMargin=20*mm,
        leftMargin=20*mm, rightMargin=20*mm,
    )
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Title"],
        fontSize=18, spaceAfter=6, textColor=colors.HexColor("#1e293b"),
    )
    h2_style = ParagraphStyle(
        "SecHeading", parent=styles["Heading2"],
        fontSize=12, textColor=colors.HexColor("#334155"),
        spaceBefore=12, spaceAfter=4,
    )
    label_style = ParagraphStyle(
        "FieldLabel", parent=styles["Normal"],
        fontSize=9, textColor=colors.HexColor("#64748b"),
        spaceAfter=1,
    )
    value_style = ParagraphStyle(
        "FieldValue", parent=styles["Normal"],
        fontSize=10, textColor=colors.HexColor("#0f172a"),
        spaceAfter=8,
    )
    h_field_style = ParagraphStyle(
        "FieldHeading", parent=styles["Heading3"],
        fontSize=11, textColor=colors.HexColor("#1e293b"),
        spaceBefore=8, spaceAfter=4,
    )

    story = []
    story.append(Paragraph(template.name, title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#e2e8f0")))
    story.append(Spacer(1, 8))

    for section in template.sections:
        story.append(Paragraph(section["title"], h2_style))
        if section.get("description"):
            story.append(Paragraph(section["description"], label_style))
        story.append(Spacer(1, 4))

        for field in section.get("fields", []):
            ftype = field.get("type", "text")
            key = field.get("key", "")
            label = field.get("label", "")
            value = _display_value(values.get(key, ""))

            if ftype == "divider":
                story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#e2e8f0")))
                story.append(Spacer(1, 4))
                continue

            if ftype == "heading":
                story.append(Paragraph(label, h_field_style))
                continue

            if ftype == "boolean":
                checked = "☑" if value else "☐"
                story.append(Paragraph(f"{checked}  {label}", value_style))
                continue

            if ftype == "table":
                rows = value if isinstance(value, list) and value else []
                cols = field.get("columns", [])
                if rows and cols:
                    header = [col["label"] for col in cols]
                    tdata = [header] + [
                        [str(row.get(col["key"], "")) for col in cols]
                        for row in rows
                    ]
                    col_width = (doc.width) / len(cols)
                    t = Table(tdata, colWidths=[col_width]*len(cols), repeatRows=1)
                    t.setStyle(TableStyle([
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f1f5f9")),
                        ("TEXTCOLOR",  (0, 0), (-1, 0), colors.HexColor("#334155")),
                        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("FONTSIZE",   (0, 0), (-1, 0), 9),
                        ("FONTSIZE",   (0, 1), (-1, -1), 9),
                        ("GRID",       (0, 0), (-1, -1), 0.5, colors.HexColor("#e2e8f0")),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
                        ("TOPPADDING",  (0, 0), (-1, -1), 5),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                        ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ]))
                    story.append(Paragraph(label, label_style))
                    story.append(t)
                    story.append(Spacer(1, 6))
                continue

            if ftype == "signature":
                story.append(Paragraph(label, label_style))
                raw = _decode_data_url_image(value)
                if raw:
                    from reportlab.platypus import Image as RLImage
                    try:
                        img = RLImage(BytesIO(raw))
                        iw, ih = (img.imageWidth or 1), (img.imageHeight or 1)
                        w = min(60 * mm, doc.width)
                        img.drawWidth, img.drawHeight = w, w * ih / iw
                        story.append(img)
                    except Exception:
                        story.append(Paragraph("[signature]", value_style))
                else:
                    story.append(Paragraph("—", value_style))
                story.append(Spacer(1, 6))
                continue

            # Default field
            story.append(Paragraph(label, label_style))
            story.append(Paragraph(str(value) if value else "—", value_style))

    doc.build(story)
    return buf.getvalue()


def generate_built_docx(template, values) -> bytes:
    """Generate DOCX from built template using python-docx."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Set margins
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # Title
    title_para = doc.add_heading(template.name, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

    for tmpl_section in template.sections:
        doc.add_heading(tmpl_section["title"], 1)
        if tmpl_section.get("description"):
            desc_para = doc.add_paragraph(tmpl_section["description"])
            desc_para.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8b)
            desc_para.runs[0].font.size = Pt(9)

        for field in tmpl_section.get("fields", []):
            ftype = field.get("type", "text")
            key = field.get("key", "")
            label = field.get("label", "")
            value = _display_value(values.get(key, ""))

            if ftype == "divider":
                p = doc.add_paragraph()
                pPr = p._p.get_or_add_pPr()
                pBdr = OxmlElement("w:pBdr")
                bottom = OxmlElement("w:bottom")
                bottom.set(qn("w:val"), "single")
                bottom.set(qn("w:sz"), "6")
                bottom.set(qn("w:space"), "1")
                bottom.set(qn("w:color"), "E2E8F0")
                pBdr.append(bottom)
                pPr.append(pBdr)
                continue

            if ftype == "heading":
                doc.add_heading(label, 3)
                continue

            if ftype == "boolean":
                p = doc.add_paragraph()
                p.add_run("☑ " if value else "☐ ").bold = True
                p.add_run(label)
                continue

            if ftype == "table":
                rows = value if isinstance(value, list) and value else []
                cols = field.get("columns", [])
                if rows and cols:
                    lp = doc.add_paragraph()
                    lp.add_run(label + ":").bold = True
                    table = doc.add_table(rows=1, cols=len(cols))
                    table.style = "Table Grid"
                    hdr_cells = table.rows[0].cells
                    for i, col in enumerate(cols):
                        hdr_cells[i].text = col["label"]
                        for run in hdr_cells[i].paragraphs[0].runs:
                            run.bold = True
                    for row_data in rows:
                        row_cells = table.add_row().cells
                        for i, col in enumerate(cols):
                            row_cells[i].text = str(row_data.get(col["key"], ""))
                    doc.add_paragraph()
                continue

            if ftype == "signature":
                p = doc.add_paragraph()
                run_label = p.add_run(f"{label}: ")
                run_label.bold = True
                run_label.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
                raw = _decode_data_url_image(value)
                if raw:
                    try:
                        from docx.shared import Inches
                        doc.add_picture(BytesIO(raw), width=Inches(2))
                    except Exception:
                        doc.add_paragraph("[signature]")
                continue

            p = doc.add_paragraph()
            run_label = p.add_run(f"{label}: ")
            run_label.bold = True
            run_label.font.color.rgb = RGBColor(0x47, 0x55, 0x69)
            p.add_run(str(value) if value else "")

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── Document designer (WYSIWYG block layout) → DOCX ─────────────────────────

_TOKEN_RE = re.compile(r"\{\{([a-zA-Z0-9_]+)\}\}")
# User placeholders the admin marks for the recipient to fill when editing the
# generated document, e.g. [[Amount]]. Rendered as highlighted fill-in markers.
_PLACEHOLDER_RE = re.compile(r"\[\[([^\]]+)\]\]")


def _designer_add_text(paragraph, raw, values, *, size=None, color=None,
                       bold=None, italic=None, underline=None, font=None):
    """
    Emit runs into a paragraph for a designer text string. {{merge_fields}} are
    substituted to static text; [[user placeholders]] become highlighted runs the
    recipient fills in when editing the generated document.
    """
    from docx.shared import Pt
    from docx.enum.text import WD_COLOR_INDEX

    substituted = _subst_tokens(raw, values)
    if not substituted:
        return
    # re.split with a capture group yields: [text, label, text, label, ...] — the
    # odd indices are the placeholder labels.
    for i, part in enumerate(_PLACEHOLDER_RE.split(substituted)):
        if not part:
            continue
        run = paragraph.add_run(part)
        is_placeholder = i % 2 == 1
        if size:
            run.font.size = Pt(size)
        if font:
            run.font.name = font
        if bold is not None:
            run.bold = bold
        if italic is not None:
            run.italic = italic
        if underline is not None:
            run.underline = underline
        if is_placeholder:
            run.font.highlight_color = WD_COLOR_INDEX.YELLOW
        elif color is not None:
            run.font.color.rgb = color


def _subst_tokens(text, values, keep_page=False):
    """Replace {{key}} with values. When keep_page is True, leave {{page}} and
    {{pages}} intact so they can become Word page-number fields."""
    if not text:
        return ""

    def repl(m):
        key = m.group(1)
        if keep_page and key in ("page", "pages"):
            return m.group(0)
        val = values.get(key)
        return "" if val is None else str(val)

    return _TOKEN_RE.sub(repl, str(text))


def _designer_hex_to_rgb(hexstr, default="1F2933"):
    from docx.shared import RGBColor
    s = (hexstr or "").lstrip("#")
    if len(s) != 6:
        s = default
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return RGBColor(0x1F, 0x29, 0x33)


def _designer_font_family(css):
    if not css:
        return None
    first = css.split(",")[0].strip().strip("'\"")
    return first or None


def _designer_add_page_field(paragraph, instr):
    """Append a Word field (PAGE / NUMPAGES) to a paragraph."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), instr)
    run = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "1"
    run.append(t)
    fld.append(run)
    paragraph._p.append(fld)


def _designer_band_runs(paragraph, text, values):
    """Render header/footer band text into a paragraph, turning {{page}} /
    {{pages}} into live Word fields and substituting other tokens."""
    rendered = _subst_tokens(text, values, keep_page=True)
    for part in re.split(r"(\{\{page\}\}|\{\{pages\}\})", rendered):
        if part == "{{page}}":
            _designer_add_page_field(paragraph, "PAGE")
        elif part == "{{pages}}":
            _designer_add_page_field(paragraph, "NUMPAGES")
        elif part:
            paragraph.add_run(part)


def _designer_clear_table_borders(table):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    tblPr = table._tbl.tblPr
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "none")
        borders.append(el)
    tblPr.append(borders)


def _designer_merge_values(values, *, user, reference_number, title):
    """
    Resolve a designer ("document") template's merge fields. These templates are
    NOT filled by the user with placeholders — their {{tokens}} auto-populate
    from the acting user, the current date, the assigned reference, the document
    title, and any document-type metadata the user supplied at upload.

    User-supplied values (the document-type metadata) take precedence; the
    context defaults below only fill tokens that were left blank. Advanced
    sources (per-field formulas, document references) are a later enhancement.
    """
    from django.utils import timezone
    from apps.documents.form_formulas import resolve_formula

    resolved = dict(values or {})
    now = timezone.localtime()

    # Canonical formula vocabulary — shared with the form builder so the same
    # formula keys resolve identically in both. The designer's "formula picker"
    # inserts these keys directly.
    formula = {
        key: resolve_formula(key, user=user, reference_number=reference_number, now=now)
        for key in (
            "current_user", "current_user_email", "current_user_department",
            "today", "now", "reference_number",
        )
    }

    # Organization identity from DMS settings (auto-fills company merge fields).
    org_name = org_address = ""
    try:
        from apps.documents.models import DMSSettings
        settings_row = DMSSettings.objects.first()
        if settings_row:
            org_name = settings_row.organization_name or ""
            org_address = settings_row.organization_address or ""
    except Exception:
        pass

    defaults = {
        # Canonical formula keys.
        **formula,
        # Friendly aliases (and designer-specific fields) → same resolved values.
        "author_name": formula["current_user"],
        "author_email": formula["current_user_email"],
        "user_name": formula["current_user"],
        "prepared_by": formula["current_user"],
        "department": formula["current_user_department"],
        "document_date": formula["today"],
        "date": formula["today"],
        "document_no": reference_number,
        "document_number": reference_number,
        "company_name": org_name,
        "company_address": org_address,
        "organization_name": org_name,
        "document_title": title,
        "title": title,
    }
    for key, value in defaults.items():
        current = resolved.get(key)
        if value and (current is None or (isinstance(current, str) and not current.strip())):
            resolved[key] = value

    # ── Document references ───────────────────────────────────────────────────
    # A value picked from a related document ({id,label,source}) resolves to its
    # label for {{key}}, and pulls common fields for {{key__field}} — e.g.
    # {{related_po}}, {{related_po__supplier}}, {{related_po__reference_number}}.
    from apps.documents.form_attachments import is_reference_value
    from apps.documents.models import Document as _RefDoc
    for key, val in list(resolved.items()):
        if not is_reference_value(val):
            continue
        resolved[key] = val.get("label", "") or ""
        ref_id = val.get("id")
        if not ref_id:
            continue
        try:
            ref_doc = _RefDoc.objects.filter(pk=ref_id).first()
        except Exception:
            ref_doc = None
        if not ref_doc:
            continue
        # Authoritative label from the referenced document (don't trust the client).
        resolved[key] = (
            f"{ref_doc.title} ({ref_doc.reference_number})"
            if ref_doc.reference_number else (ref_doc.title or "")
        )
        amt = getattr(ref_doc, "amount", None)
        dd = getattr(ref_doc, "document_date", None)
        resolved[f"{key}__reference_number"] = ref_doc.reference_number or ""
        resolved[f"{key}__title"] = ref_doc.title or ""
        resolved[f"{key}__supplier"] = getattr(ref_doc, "supplier", "") or ""
        resolved[f"{key}__amount"] = str(amt) if amt is not None else ""
        resolved[f"{key}__currency"] = getattr(ref_doc, "currency", "") or ""
        resolved[f"{key}__document_date"] = dd.strftime("%d %b %Y") if dd else ""
        for mk, mv in (getattr(ref_doc, "metadata", None) or {}).items():
            if isinstance(mv, (str, int, float)):
                resolved.setdefault(f"{key}__{mk}", str(mv))

    return resolved


def generate_designer_docx(design, values) -> bytes:
    """
    Render a WYSIWYG document-designer template (block layout) to an editable
    DOCX, substituting {{tokens}} from `values`. Faithful-but-approximate to the
    on-screen designer; the output follows the normal Office editing lifecycle.
    """
    from docx import Document
    from docx.shared import Pt, Mm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.section import WD_ORIENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    design = design or {}
    theme = design.get("theme", {}) or {}
    page = design.get("page", {}) or {}
    blocks = design.get("blocks", []) or []

    body_font = _designer_font_family(theme.get("fontFamily")) or "Calibri"
    heading_font = _designer_font_family(theme.get("headingFamily")) or body_font
    base_px = theme.get("baseFontSize") or 13
    base_pt = max(8, round(float(base_px) * 0.75))
    text_rgb = _designer_hex_to_rgb(theme.get("textColor"), "1F2933")
    heading_rgb = _designer_hex_to_rgb(theme.get("headingColor"), "0F2A3A")
    accent_rgb = _designer_hex_to_rgb(theme.get("accentColor"), "287EAD")

    px_to_pt = lambda px, fallback: max(8, round(float(px) * 0.75)) if px else fallback

    doc = Document()

    # Base style
    normal = doc.styles["Normal"]
    normal.font.name = body_font
    normal.font.size = Pt(base_pt)
    normal.font.color.rgb = text_rgb

    ALIGN = {
        "left": WD_ALIGN_PARAGRAPH.LEFT, "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT, "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }

    # ── Page setup ──────────────────────────────────────────────────────────
    PAGE_DIMS = {"A4": (210, 297), "Letter": (216, 279), "Legal": (216, 356)}
    w_mm, h_mm = PAGE_DIMS.get(page.get("size", "A4"), (210, 297))
    margin = page.get("margin", {}) or {}
    for section in doc.sections:
        if page.get("orientation") == "landscape":
            section.orientation = WD_ORIENT.LANDSCAPE
            section.page_width, section.page_height = Mm(h_mm), Mm(w_mm)
        else:
            section.page_width, section.page_height = Mm(w_mm), Mm(h_mm)
        section.top_margin = Mm(margin.get("top", 20))
        section.bottom_margin = Mm(margin.get("bottom", 20))
        section.left_margin = Mm(margin.get("left", 18))
        section.right_margin = Mm(margin.get("right", 18))

    # ── Header / footer bands (left / center / right) ───────────────────────
    def render_band(band, container):
        if not band or not band.get("enabled"):
            return
        content = band.get("content", {}) or {}
        table = container.add_table(rows=1, cols=3, width=Mm(w_mm - margin.get("left", 18) - margin.get("right", 18)))
        _designer_clear_table_borders(table)
        cells = table.rows[0].cells
        for cell, (slot, align) in zip(cells, (("left", "left"), ("center", "center"), ("right", "right"))):
            para = cell.paragraphs[0]
            para.alignment = ALIGN[align]
            _designer_band_runs(para, content.get(slot, ""), values)

    section = doc.sections[0]
    render_band(design.get("header"), section.header)
    render_band(design.get("footer"), section.footer)

    # ── Blocks ──────────────────────────────────────────────────────────────
    def add_text(paragraph, raw, **style):
        _designer_add_text(paragraph, raw, values, **style)

    for b in blocks:
        btype = b.get("type")
        align = ALIGN.get(b.get("align", "left"), WD_ALIGN_PARAGRAPH.LEFT)

        if btype == "heading":
            level = b.get("level", 2)
            size = px_to_pt(b.get("fontSize"), {1: 18, 2: 14, 3: 12}.get(level, 14))
            p = doc.add_paragraph()
            p.alignment = align
            add_text(p, b.get("text", ""), size=size, font=heading_font, bold=True,
                     color=_designer_hex_to_rgb(b.get("color")) if b.get("color") else heading_rgb)

        elif btype in ("paragraph", "quote"):
            p = doc.add_paragraph()
            p.alignment = align
            add_text(p, b.get("text", ""), size=px_to_pt(b.get("fontSize"), None),
                     color=_designer_hex_to_rgb(b.get("color")) if b.get("color") else None,
                     bold=b.get("bold"), italic=b.get("italic") or (btype == "quote"),
                     underline=b.get("underline"))

        elif btype in ("bulleted_list", "numbered_list"):
            style = "List Bullet" if btype == "bulleted_list" else "List Number"
            for item in b.get("items", []) or []:
                p = doc.add_paragraph(style=style)
                add_text(p, item)

        elif btype == "key_value":
            pairs = b.get("pairs", []) or []
            if pairs:
                table = doc.add_table(rows=len(pairs), cols=2)
                _designer_clear_table_borders(table)
                for i, pair in enumerate(pairs):
                    lc, vc = table.rows[i].cells
                    add_text(lc.paragraphs[0], pair.get("label", ""), bold=True)
                    add_text(vc.paragraphs[0], pair.get("value", ""))

        elif btype == "data_table":
            cols = b.get("columns", []) or []
            if cols:
                if b.get("bound"):
                    # Bound to a collection if one was supplied (one row per
                    # record); otherwise emit blank fillable rows for the user to
                    # complete when editing the generated document.
                    source = values.get(b.get("sourceKey") or "")
                    if isinstance(source, list) and source:
                        rows = [
                            [str((rec or {}).get(c.get("key", ""), "")) for c in cols]
                            for rec in source
                        ]
                    else:
                        n = max(1, int(b.get("fillRows") or 3))
                        rows = [["" for _ in cols] for _ in range(n)]
                else:
                    rows = b.get("rows", []) or []
                table = doc.add_table(rows=1, cols=len(cols))
                table.style = "Table Grid" if b.get("bordered", True) else "Light List"
                for i, col in enumerate(cols):
                    cell = table.rows[0].cells[i]
                    run = cell.paragraphs[0].add_run(col.get("label", ""))
                    run.bold = True
                for row in rows:
                    cells = table.add_row().cells
                    for i, _col in enumerate(cols):
                        val = row[i] if i < len(row) else ""
                        add_text(cells[i].paragraphs[0], val)

        elif btype == "two_column":
            table = doc.add_table(rows=1, cols=2)
            _designer_clear_table_borders(table)
            lc, rc = table.rows[0].cells
            add_text(lc.paragraphs[0], b.get("left", ""))
            add_text(rc.paragraphs[0], b.get("right", ""))

        elif btype == "divider":
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "C8CDD2")
            pBdr.append(bottom)
            pPr.append(pBdr)

        elif btype == "spacer":
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(px_to_pt(b.get("height", 24), 18))

        elif btype == "page_break":
            doc.add_page_break()

        elif btype == "signature":
            sigs = b.get("signatories", []) or []
            if sigs:
                # Lay out up to 3 signatories per row, wrapping into more rows for
                # additional approvers. Each cell: role, a signing line, then the
                # Name/Date fields (which may be {{auto-fill}}, [[placeholders]] or
                # plain text).
                per_row = 3
                for start in range(0, len(sigs), per_row):
                    chunk = sigs[start:start + per_row]
                    table = doc.add_table(rows=1, cols=len(chunk))
                    _designer_clear_table_borders(table)
                    for i, s in enumerate(chunk):
                        cell = table.rows[0].cells[i]
                        cell.paragraphs[0].add_run(s.get("role", "")).bold = True
                        cell.add_paragraph()  # signing space
                        cell.add_paragraph().add_run("_______________________")
                        lbl = cell.add_paragraph().add_run("Signature")
                        lbl.italic = True
                        lbl.font.size = Pt(8)
                        if s.get("nameToken"):
                            p = cell.add_paragraph()
                            p.add_run("Name: ").bold = True
                            add_text(p, s.get("nameToken"))
                        if s.get("dateToken"):
                            p = cell.add_paragraph()
                            p.add_run("Date: ").bold = True
                            add_text(p, s.get("dateToken"))
                    doc.add_paragraph()  # spacing between signatory rows

        elif btype in ("image", "logo"):
            # Admin-uploaded images arrive as base64 data URLs embedded in the
            # template; decode and embed them. Fall back to a light text marker
            # for a missing/URL-only source.
            from docx.shared import Emu
            raw_img = _decode_data_url_image(b.get("src"))
            width_px = b.get("width") or 160
            emu_width = Emu(int(float(width_px) / 96 * 914400))  # 96 px/inch

            def _place_image(paragraph):
                if raw_img:
                    try:
                        paragraph.add_run().add_picture(BytesIO(raw_img), width=emu_width)
                        return
                    except Exception:
                        pass
                add_text(paragraph, _subst_tokens(b.get("alt", ""), values) or ("Logo" if btype == "logo" else "Image"), italic=True)

            side = b.get("float") or "none"
            if side in ("left", "right"):
                # Side-by-side: borderless 1×2 table — image in one cell, the
                # `beside` content in the other. This is the Word-native way to
                # place text alongside an image (docx has no CSS float/wrap).
                table = doc.add_table(rows=1, cols=2)
                _designer_clear_table_borders(table)
                img_mm = min(w_mm - margin.get("left", 18) - margin.get("right", 18) - 10,
                             max(10, float(width_px) / 96 * 25.4))
                img_cell, text_cell = table.rows[0].cells
                if side == "right":
                    img_cell, text_cell = text_cell, img_cell
                img_cell.width = Mm(img_mm)
                _place_image(img_cell.paragraphs[0])
                add_text(text_cell.paragraphs[0], b.get("beside", ""))
            else:
                p = doc.add_paragraph()
                p.alignment = align
                _place_image(p)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── Uploaded template fillers ───────────────────────────────────────────────

def fill_docx_template(template, values) -> bytes:
    """
    Fill uploaded DOCX template with placeholder values.
    Handles placeholders split across runs by stitching paragraph text.
    """
    from docx import Document

    doc = Document(template.file.path)

    def process_para(para):
        _replace_placeholder_in_paragraph(para, values)

    # Body paragraphs
    for para in doc.paragraphs:
        process_para(para)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    process_para(para)

    # Headers & footers
    for section in doc.sections:
        for hf in [section.header, section.footer]:
            if hf:
                for para in hf.paragraphs:
                    process_para(para)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def fill_xlsx_template(template, values) -> bytes:
    """Fill uploaded XLSX template with placeholder values."""
    import openpyxl
    import re

    wb = openpyxl.load_workbook(template.file.path)

    def replace_cell(cell):
        if cell.value and isinstance(cell.value, str) and "{{" in cell.value:
            def replacer(m):
                key = m.group(1)
                val = values.get(key, "")
                return str(val) if val is not None else ""
            cell.value = re.sub(r"\{\{([a-zA-Z0-9_]+)\}\}", replacer, cell.value)

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                replace_cell(cell)

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _replace_placeholder_in_pptx_paragraph(para, values: dict):
    """
    Replace {{key}} placeholders in a python-pptx paragraph, stitching text
    across runs (PowerPoint, like Word, often splits a placeholder over several
    runs). The merged text is written back into the first run; the rest cleared.
    """
    runs = para.runs
    if not runs:
        return
    full_text = "".join(run.text for run in runs)
    if "{{" not in full_text:
        return

    def replacer(match):
        val = values.get(match.group(1))
        return "" if val is None else str(val)

    new_text = re.sub(r"\{\{([a-zA-Z0-9_]+)\}\}", replacer, full_text)
    if new_text == full_text:
        return
    runs[0].text = new_text
    for run in runs[1:]:
        run.text = ""


def fill_pptx_template(template, values) -> bytes:
    """Fill an uploaded PPTX template with placeholder values."""
    from pptx import Presentation

    prs = Presentation(template.file.path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _replace_placeholder_in_pptx_paragraph(para, values)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _replace_placeholder_in_pptx_paragraph(para, values)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def pptx_to_pdf(pptx_bytes: bytes) -> bytes:
    """Convert PPTX to PDF using LibreOffice headless."""
    from apps.documents.tasks import _convert_office_source_to_pdf_bytes

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(pptx_bytes)
        tmp_path = f.name
    try:
        return _convert_office_source_to_pdf_bytes(
            Path(tmp_path), soffice_bin="libreoffice", timeout=60
        )
    finally:
        os.unlink(tmp_path)


def docx_to_pdf(docx_bytes: bytes) -> bytes:
    """Convert DOCX to PDF using LibreOffice headless."""
    from apps.documents.tasks import _convert_office_source_to_pdf_bytes

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        f.write(docx_bytes)
        tmp_path = f.name
    try:
        return _convert_office_source_to_pdf_bytes(
            Path(tmp_path), soffice_bin="libreoffice", timeout=60
        )
    finally:
        os.unlink(tmp_path)


def xlsx_to_pdf(xlsx_bytes: bytes) -> bytes:
    """Convert XLSX to PDF using LibreOffice headless."""
    from apps.documents.tasks import _convert_office_source_to_pdf_bytes

    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        f.write(xlsx_bytes)
        tmp_path = f.name
    try:
        return _convert_office_source_to_pdf_bytes(
            Path(tmp_path), soffice_bin="libreoffice", timeout=60
        )
    finally:
        os.unlink(tmp_path)


# ─── Main entry ─────────────────────────────────────────────────────────────

def generate_document_from_template_sync(template, values, fmt, title, user, type_id):
    """Generate a document from a template synchronously."""
    from apps.documents.models import Document, DocumentStatus
    from apps.documents.serializers import _generate_unique_reference
    from apps.documents.form_attachments import descriptors_to_names
    from apps.documents.form_formulas import apply_formulas
    from django.core.files.base import ContentFile
    import hashlib

    is_xlsx = template.file_name.endswith((".xlsx", ".xls")) if template.file_name else False
    is_pptx = template.file_name.endswith((".pptx", ".ppt")) if template.file_name else False
    kind = getattr(template, "kind", "form") or "form"

    # Reserve the reference up front so a `reference_number` formula can use it.
    reference_number = _generate_unique_reference(template.document_type)

    if template.type == "built" and kind == "document":
        # WYSIWYG document designer: merge fields auto-populate from the user,
        # date, reference and the document-type metadata — the user fills no
        # placeholders. Render the block layout to an editable DOCX (or PDF on
        # request) so the result follows the normal Office lifecycle.
        merge_values = _designer_merge_values(
            values, user=user, reference_number=reference_number, title=title
        )
        render_values = descriptors_to_names(merge_values)
        docx_content = generate_designer_docx(template.design, render_values)
        if fmt == "pdf":
            content = docx_to_pdf(docx_content)
            filename = f"{title}.pdf"
            content_type = "application/pdf"
        else:
            content = docx_content
            filename = f"{title}.docx"
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    elif template.type == "built":
        # Freeze auto-fill formula values authoritatively (creator, submit time,
        # assigned reference) into the stored form values.
        values = apply_formulas(
            values, template.sections, user=user, reference_number=reference_number
        )
        # Freeze calculated fields (e.g. "total_days * daily_rate") the same
        # way — computed server-side from the (now formula-frozen) values so
        # the stored document and the rendered file always agree, regardless
        # of what the client last had on screen.
        from apps.templates_engine.conditions import compute_calculated_values
        values = compute_calculated_values(template.sections, values)
        # The stored form.values keeps structured attachment descriptors and
        # reference {id,label} objects; the rendered file shows display strings.
        render_values = descriptors_to_names(values)
        if fmt == "pdf":
            content = generate_built_pdf(template, render_values)
            filename = f"{title}.pdf"
            content_type = "application/pdf"
        else:
            content = generate_built_docx(template, render_values)
            filename = f"{title}.docx"
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    else:  # uploaded
        if is_xlsx:
            xlsx_content = fill_xlsx_template(template, values)
            if fmt == "pdf":
                content = xlsx_to_pdf(xlsx_content)
                filename = f"{title}.pdf"
                content_type = "application/pdf"
            else:
                content = xlsx_content
                filename = f"{title}.xlsx"
                content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif is_pptx:
            pptx_content = fill_pptx_template(template, values)
            if fmt == "pdf":
                content = pptx_to_pdf(pptx_content)
                filename = f"{title}.pdf"
                content_type = "application/pdf"
            else:
                content = pptx_content
                filename = f"{title}.pptx"
                content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            docx_content = fill_docx_template(template, values)
            if fmt == "pdf":
                content = docx_to_pdf(docx_content)
                filename = f"{title}.pdf"
                content_type = "application/pdf"
            else:
                content = docx_content
                filename = f"{title}.docx"
                content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    checksum = hashlib.sha256(content).hexdigest()

    doc_metadata = {
        "template_id": str(template.id),
        "template_name": template.name,
        **_form_values_for_metadata(values),
    }
    if template.type == "built" and kind == "form":
        # Built FORM templates are interactive forms: store the schema snapshot +
        # the entered values so the document IS the filled form and can be
        # re-rendered / edited in-app (no external editor). The generated file is
        # just a view. Designer ("document") templates render a static editable
        # file instead and follow the normal Office lifecycle — no form snapshot.
        doc_metadata["form"] = {
            "template_id": str(template.id),
            "sections": template.sections,
            "values": values,
        }
        # Snapshot the SunSystems integration mapping (journal/budget) onto the
        # document so budget checks and journal posting depend on the document
        # alone and survive later template edits. See apps/sunsystems/config.py.
        ss_mapping = getattr(template, "sunsystems", None)
        if isinstance(ss_mapping, dict) and ss_mapping:
            doc_metadata["sunsystems"] = ss_mapping

    create_kwargs = dict(
        title=title,
        reference_number=reference_number,
        file=ContentFile(content, name=filename),
        file_name=filename,
        file_size=len(content),
        file_mime_type=content_type,
        checksum=checksum,
        uploaded_by=user,
        owned_by=user,
        document_type_id=type_id,
        is_self_upload=False,
        status=DocumentStatus.DRAFT,
        metadata=doc_metadata,
        **_document_field_kwargs(values),
        # A template-generated document is the starting point, not a user version.
        # It stays unversioned (v0 → shows as "—") until the user first edits it,
        # at which point the first save becomes version 1.
        current_version=0,
    )
    try:
        doc = Document.objects.create(**create_kwargs)
    except SEARCH_INDEX_EXCEPTIONS:
        # Elasticsearch is read-only (e.g. disk flood-stage). The row is already
        # committed; fetch it so document creation still succeeds. Indexing will
        # catch up once ES recovers.
        logger.warning(
            "Template document %s saved but realtime indexing failed (ES read-only).",
            reference_number,
        )
        doc = Document.objects.get(reference_number=reference_number)

    if doc.is_office_doc():
        try:
            from apps.documents.tasks import generate_document_preview

            Document.objects.filter(id=doc.id, preview_status="").update(
                preview_status="pending"
            )
            generate_document_preview.delay(str(doc.id))
        except Exception:
            pass
    try:
        from apps.search.indexing import schedule_document_search_pipeline

        schedule_document_search_pipeline(
            str(doc.id),
            reextract_content=True,
            index_immediately=True,
        )
    except Exception:
        pass
    return doc